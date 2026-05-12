"""DocuVault Desktop — Local Python backend.
Runs as a Tauri sidecar. SQLite database, no Docker, files indexed in-place.
"""

import hashlib
import io
import json
import os
import sqlite3
import sys
import uuid
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path

import magic
import pytesseract
from fastapi import FastAPI, Query, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from PIL import Image
from pdf2image import convert_from_bytes
from PyPDF2 import PdfReader
import anthropic
import meilisearch
import uvicorn

# --- Config ---

APP_DATA = Path.home() / ".docuvault"
APP_DATA.mkdir(exist_ok=True)
DB_PATH = APP_DATA / "docuvault.db"
MEILI_URL = "http://127.0.0.1:7701"
MEILI_KEY = "docuvault-local"

SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf",
    ".doc", ".docx", ".dot", ".dotx",
    ".xls", ".xlsx", ".xlt", ".xltx",
    ".ppt", ".pptx", ".pot", ".potx",
    ".odt", ".ods", ".odp", ".odg", ".odf",  # LibreOffice
    ".rtf",
    ".txt", ".csv", ".tsv", ".md", ".log", ".json", ".xml", ".yaml", ".yml",
    ".html", ".htm", ".xhtml", ".mhtml",
    ".tex", ".latex",
    ".epub", ".mobi",
    # Images (OCR-able)
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp",
    ".gif", ".heic", ".heif", ".avif",
    ".svg",
    ".raw", ".cr2", ".nef", ".arw", ".dng", ".rw2",  # RAW photos
    ".ico",
    # Scanned documents
    ".djvu", ".djv",
    # eBooks
    ".fb2",
    # Email
    ".eml", ".msg",
    # Archives (index metadata, list contents)
    ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz",
    # CAD / Engineering
    ".dwg", ".dxf", ".sldprt", ".sldasm", ".slddrw",
    ".step", ".stp", ".iges", ".igs", ".stl", ".3mf", ".obj",
    ".f3d", ".ipt", ".iam",
    # Code
    ".py", ".js", ".ts", ".tsx", ".jsx",
    ".java", ".c", ".cpp", ".h", ".hpp",
    ".go", ".rs", ".rb", ".php", ".swift", ".kt",
    ".css", ".scss", ".less",
    ".sql", ".sh", ".bash", ".zsh", ".ps1",
    ".r", ".m", ".ipynb",
    # Data
    ".parquet", ".avro", ".feather",
    # Design
    ".psd", ".ai", ".eps", ".sketch", ".fig", ".xd",
    ".indd",
    # Audio (metadata only)
    ".mp3", ".m4a", ".wav", ".flac", ".ogg", ".aac", ".wma", ".opus",
    # Video (metadata only)
    ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm",
    ".m4v", ".mpg", ".mpeg", ".3gp", ".mts", ".m2ts",
    # Misc
    ".ics", ".vcf", ".vcard",  # Calendar, contacts
    ".bib",  # Bibliography
    ".pages", ".numbers", ".key",  # Apple iWork
}

# --- Database ---


def get_db():
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA foreign_keys=ON")
    return conn


def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            file_path TEXT NOT NULL UNIQUE,
            title TEXT NOT NULL,
            original_filename TEXT,
            category TEXT,
            subcategory TEXT,
            file_hash TEXT,
            mime_type TEXT,
            file_size INTEGER,
            page_count INTEGER,
            language TEXT DEFAULT 'en',
            summary TEXT,
            full_text TEXT,
            ai_confidence REAL,
            document_date TEXT,
            tags TEXT DEFAULT '[]',
            people TEXT DEFAULT '[]',
            organizations TEXT DEFAULT '[]',
            ai_metadata TEXT,
            is_archived INTEGER DEFAULT 0,
            is_deleted INTEGER DEFAULT 0,
            source TEXT,
            created_at TEXT,
            indexed_at TEXT
        );

        CREATE TABLE IF NOT EXISTS watched_folders (
            id TEXT PRIMARY KEY,
            path TEXT NOT NULL UNIQUE,
            last_scan TEXT,
            file_count INTEGER DEFAULT 0,
            is_active INTEGER DEFAULT 1
        );

        CREATE INDEX IF NOT EXISTS idx_docs_category ON documents(category);
        CREATE INDEX IF NOT EXISTS idx_docs_hash ON documents(file_hash);
        CREATE INDEX IF NOT EXISTS idx_docs_path ON documents(file_path);
        CREATE INDEX IF NOT EXISTS idx_docs_deleted ON documents(is_deleted);

        -- Users & Auth (RBAC, MFA)
        CREATE TABLE IF NOT EXISTS users (
            id TEXT PRIMARY KEY,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            name TEXT,
            role TEXT DEFAULT 'editor',
            is_active INTEGER DEFAULT 1,
            mfa_secret TEXT,
            mfa_enabled INTEGER DEFAULT 0,
            created_at TEXT
        );

        -- Audit Logs
        CREATE TABLE IF NOT EXISTS audit_log (
            id TEXT PRIMARY KEY,
            user_id TEXT,
            action TEXT NOT NULL,
            resource_type TEXT,
            resource_id TEXT,
            details TEXT,
            ip_address TEXT,
            timestamp TEXT DEFAULT (datetime('now'))
        );
        CREATE INDEX IF NOT EXISTS idx_audit_timestamp ON audit_log(timestamp);
        CREATE INDEX IF NOT EXISTS idx_audit_user ON audit_log(user_id);

        -- Version Control
        CREATE TABLE IF NOT EXISTS document_versions (
            id TEXT PRIMARY KEY,
            document_id TEXT NOT NULL,
            version_number INTEGER NOT NULL,
            file_hash TEXT,
            file_size INTEGER,
            changed_by TEXT,
            change_summary TEXT,
            snapshot_path TEXT,
            created_at TEXT DEFAULT (datetime('now')),
            FOREIGN KEY (document_id) REFERENCES documents(id)
        );

        -- Cross-links between documents
        CREATE TABLE IF NOT EXISTS document_links (
            id TEXT PRIMARY KEY,
            source_id TEXT NOT NULL,
            target_id TEXT NOT NULL,
            link_type TEXT DEFAULT 'related',
            created_at TEXT DEFAULT (datetime('now')),
            FOREIGN KEY (source_id) REFERENCES documents(id),
            FOREIGN KEY (target_id) REFERENCES documents(id)
        );

        -- Shared links
        CREATE TABLE IF NOT EXISTS shared_links (
            id TEXT PRIMARY KEY,
            document_id TEXT NOT NULL,
            token TEXT UNIQUE NOT NULL,
            expires_at TEXT,
            password_hash TEXT,
            access_count INTEGER DEFAULT 0,
            max_access INTEGER,
            created_by TEXT,
            created_at TEXT DEFAULT (datetime('now')),
            FOREIGN KEY (document_id) REFERENCES documents(id)
        );

        -- Retention policies
        CREATE TABLE IF NOT EXISTS retention_policies (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            category_match TEXT,
            subcategory_match TEXT,
            retention_days INTEGER,
            action TEXT DEFAULT 'archive',
            is_active INTEGER DEFAULT 1,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Folder templates
        CREATE TABLE IF NOT EXISTS folder_templates (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            structure TEXT NOT NULL,
            description TEXT,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Webhooks (Zapier/Make integration)
        CREATE TABLE IF NOT EXISTS webhooks (
            id TEXT PRIMARY KEY,
            url TEXT NOT NULL,
            events TEXT NOT NULL,
            secret TEXT,
            is_active INTEGER DEFAULT 1,
            last_triggered TEXT,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Integration configs
        CREATE TABLE IF NOT EXISTS integrations (
            id TEXT PRIMARY KEY,
            provider TEXT NOT NULL,
            config TEXT NOT NULL,
            is_active INTEGER DEFAULT 1,
            last_sync TEXT,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Comments on documents
        CREATE TABLE IF NOT EXISTS comments (
            id TEXT PRIMARY KEY,
            document_id TEXT NOT NULL,
            parent_id TEXT,
            user_id TEXT,
            user_name TEXT DEFAULT 'You',
            content TEXT NOT NULL,
            created_at TEXT DEFAULT (datetime('now')),
            updated_at TEXT,
            FOREIGN KEY (document_id) REFERENCES documents(id),
            FOREIGN KEY (parent_id) REFERENCES comments(id)
        );
        CREATE INDEX IF NOT EXISTS idx_comments_doc ON comments(document_id);

        -- Smart folders (saved searches)
        CREATE TABLE IF NOT EXISTS smart_folders (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            icon TEXT DEFAULT 'folder-search',
            filters TEXT NOT NULL,
            sort_by TEXT DEFAULT 'indexed_at',
            sort_order TEXT DEFAULT 'desc',
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Activity log (user-facing)
        CREATE TABLE IF NOT EXISTS activity_log (
            id TEXT PRIMARY KEY,
            action TEXT NOT NULL,
            description TEXT NOT NULL,
            document_id TEXT,
            document_title TEXT,
            metadata TEXT,
            created_at TEXT DEFAULT (datetime('now'))
        );
        CREATE INDEX IF NOT EXISTS idx_activity_created ON activity_log(created_at);

        -- Workspaces (multi-tenant)
        CREATE TABLE IF NOT EXISTS workspaces (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            owner_id TEXT,
            description TEXT,
            is_active INTEGER DEFAULT 1,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Document permissions (per-doc/folder)
        CREATE TABLE IF NOT EXISTS document_permissions (
            id TEXT PRIMARY KEY,
            document_id TEXT,
            category_match TEXT,
            user_id TEXT,
            permission TEXT DEFAULT 'read',
            granted_by TEXT,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- Document templates (AI-generated)
        CREATE TABLE IF NOT EXISTS document_templates (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            description TEXT,
            template_type TEXT NOT NULL,
            content TEXT NOT NULL,
            variables TEXT,
            created_at TEXT DEFAULT (datetime('now'))
        );

        -- App settings (key-value)
        CREATE TABLE IF NOT EXISTS app_settings (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL,
            updated_at TEXT DEFAULT (datetime('now'))
        );

        -- Keyboard shortcuts
        CREATE TABLE IF NOT EXISTS keyboard_shortcuts (
            id TEXT PRIMARY KEY,
            key_combo TEXT NOT NULL UNIQUE,
            action TEXT NOT NULL,
            description TEXT,
            is_active INTEGER DEFAULT 1
        );

        -- Plugin registry
        CREATE TABLE IF NOT EXISTS plugins (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            version TEXT,
            description TEXT,
            entry_point TEXT,
            config TEXT,
            is_active INTEGER DEFAULT 1,
            installed_at TEXT DEFAULT (datetime('now'))
        );
    """)
    conn.commit()
    conn.close()


# --- Text Extraction ---
#
# Strategy:
#   1. Try `unstructured` first — best quality for 26 core doc types
#      (PDF, DOCX, PPTX, XLSX, HTML, EML, EPUB, images with OCR, etc.)
#   2. Fall back to custom extractors for types unstructured doesn't handle
#      (audio/video metadata, archives, CAD, code, Jupyter, etc.)

# Types that `unstructured` handles well
UNSTRUCTURED_TYPES = {
    ".bmp", ".csv", ".doc", ".docx", ".eml", ".epub", ".heic",
    ".html", ".htm", ".jpeg", ".jpg", ".png", ".md", ".msg",
    ".odt", ".org", ".p7s", ".pdf", ".ppt", ".pptx", ".rst",
    ".rtf", ".tiff", ".tif", ".txt", ".tsv", ".xls", ".xlsx", ".xml",
}


def extract_text(file_path: str, mime_type: str) -> str:
    """Extract text from any supported file type."""
    path = Path(file_path)
    ext = path.suffix.lower()

    # Tier 1: Use unstructured for supported types (best quality)
    if ext in UNSTRUCTURED_TYPES:
        text = _extract_with_unstructured(path)
        if text and len(text.strip()) > 10:
            return text

    # Tier 2: Custom extractors for everything else
    # Code / plain text / config / data
    if ext in {".log", ".json", ".yaml", ".yml", ".py", ".js", ".ts", ".tsx",
               ".jsx", ".java", ".c", ".cpp", ".h", ".hpp", ".go", ".rs",
               ".rb", ".php", ".swift", ".kt", ".css", ".scss", ".less",
               ".sql", ".sh", ".bash", ".zsh", ".ps1", ".r", ".m",
               ".tex", ".latex", ".bib", ".ics", ".vcf", ".vcard",
               ".xhtml", ".mhtml", ".svg"}:
        try:
            return path.read_text(errors="replace")[:50000]
        except Exception:
            return ""

    # Images not handled by unstructured (RAW photos)
    if ext in {".webp", ".gif", ".avif", ".raw", ".cr2", ".nef",
               ".arw", ".dng", ".rw2", ".ico"}:
        return _ocr_image(path)

    # Jupyter notebooks
    if ext == ".ipynb":
        return _extract_ipynb(path)

    # Archives — list contents
    if ext in {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz"}:
        return _extract_archive_listing(path)

    # Audio/Video — extract metadata
    if ext in {".mp3", ".m4a", ".wav", ".flac", ".ogg", ".aac", ".wma", ".opus",
               ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm",
               ".m4v", ".mpg", ".mpeg", ".3gp", ".mts", ".m2ts"}:
        return _extract_media_metadata(path)

    # LibreOffice formats not in unstructured
    if ext in {".ods", ".odp", ".odg", ".odf"}:
        return _extract_odf(path)

    # Presentations not in unstructured
    if ext in {".potx", ".dotx", ".xltx"}:
        return _extract_pptx(path) if "pot" in ext else ""

    # CAD / Design — record as metadata
    if ext in {".dwg", ".dxf", ".sldprt", ".sldasm", ".slddrw", ".step", ".stp",
               ".iges", ".igs", ".stl", ".3mf", ".obj", ".f3d", ".ipt", ".iam",
               ".psd", ".ai", ".eps", ".sketch", ".fig", ".xd", ".indd",
               ".pages", ".numbers", ".key"}:
        return f"[{ext.upper().lstrip('.')} file] {path.name}"

    # Apple iWork
    if ext in {".pages", ".numbers", ".key"}:
        return f"[Apple iWork] {path.name}"

    # Fallback — try reading as text
    try:
        return path.read_text(errors="replace")[:20000]
    except Exception:
        return f"[Binary file] {path.name}"


def _extract_with_unstructured(path: Path) -> str:
    """Use the unstructured library for high-quality extraction with layout analysis."""
    try:
        from unstructured.partition.auto import partition
        elements = partition(filename=str(path))
        texts = [str(el) for el in elements if str(el).strip()]
        return "\n\n".join(texts)[:50000]
    except Exception as e:
        # Fall back to our basic extractors
        ext = path.suffix.lower()
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext in {".docx", ".dotx"}:
            return _extract_docx(path)
        elif ext in {".xlsx", ".xltx"}:
            return _extract_xlsx(path)
        elif ext in {".pptx"}:
            return _extract_pptx(path)
        elif ext in {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".heic"}:
            return _ocr_image(path)
        elif ext in {".html", ".htm"}:
            return _extract_html(path)
        elif ext in {".eml"}:
            return _extract_eml(path)
        elif ext in {".txt", ".csv", ".tsv", ".md", ".rst", ".org", ".xml"}:
            return path.read_text(errors="replace")[:50000]
        return ""


def _extract_pdf(path: Path) -> str:
    data = path.read_bytes()
    try:
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages[:20]:
            text = page.extract_text()
            if text:
                pages.append(text)
        text = "\n\n".join(pages)
        if text and len(text.strip()) > 50:
            return text.strip()
    except Exception:
        pass
    try:
        images = convert_from_bytes(data, first_page=1, last_page=10, dpi=200)
        return "\n\n".join(
            pytesseract.image_to_string(img, lang="eng").strip()
            for img in images
        )
    except Exception:
        return ""


def _ocr_image(path: Path) -> str:
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img, lang="eng").strip()
    except Exception:
        return ""


def _extract_docx(path: Path) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        return ""


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        lines = []
        for sheet in wb.sheetnames:
            for row in wb[sheet].iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)[:50000]
    except Exception:
        return ""


def _extract_doc_legacy(path: Path) -> str:
    """Extract text from legacy .doc files."""
    try:
        import subprocess
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout[:50000]
    except Exception:
        pass
    try:
        # Fallback: try antiword or strings
        import subprocess
        result = subprocess.run(["strings", str(path)], capture_output=True, text=True, timeout=10)
        return result.stdout[:20000]
    except Exception:
        return ""


def _extract_xls_legacy(path: Path) -> str:
    """Extract text from legacy .xls files."""
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        lines = []
        for sheet in wb.sheets():
            for row_idx in range(min(sheet.nrows, 500)):
                cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)
                         if sheet.cell_value(row_idx, col)]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)[:50000]
    except Exception:
        return ""


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n\n".join(texts)[:50000]
    except Exception:
        return ""


def _extract_odf(path: Path) -> str:
    """Extract text from LibreOffice ODF files (odt, ods, odp)."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(str(path)) as zf:
            if "content.xml" in zf.namelist():
                content = zf.read("content.xml")
                root = ET.fromstring(content)
                texts = [elem.text for elem in root.iter() if elem.text and elem.text.strip()]
                return "\n".join(texts)[:50000]
    except Exception:
        return ""


def _extract_rtf(path: Path) -> str:
    """Extract text from RTF files."""
    try:
        import subprocess
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout[:50000]
    except Exception:
        pass
    try:
        content = path.read_text(errors="replace")
        import re
        text = re.sub(r'\\[a-z]+\d*\s?', '', content)
        text = re.sub(r'[{}]', '', text)
        return text[:50000]
    except Exception:
        return ""


def _extract_html(path: Path) -> str:
    """Extract text from HTML files."""
    try:
        from html.parser import HTMLParser
        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.texts = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip and data.strip():
                    self.texts.append(data.strip())

        parser = TextExtractor()
        parser.feed(path.read_text(errors="replace"))
        return "\n".join(parser.texts)[:50000]
    except Exception:
        return ""


def _extract_eml(path: Path) -> str:
    """Extract text from email .eml files."""
    try:
        import email
        msg = email.message_from_bytes(path.read_bytes())
        parts = []
        subject = msg.get("Subject", "")
        sender = msg.get("From", "")
        date = msg.get("Date", "")
        parts.append(f"Subject: {subject}\nFrom: {sender}\nDate: {date}\n")
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode(errors="replace"))
        return "\n\n".join(parts)[:50000]
    except Exception:
        return ""


def _extract_msg(path: Path) -> str:
    """Extract text from Outlook .msg files."""
    try:
        import subprocess
        result = subprocess.run(["strings", str(path)], capture_output=True, text=True, timeout=10)
        return result.stdout[:20000]
    except Exception:
        return ""


def _extract_epub(path: Path) -> str:
    """Extract text from EPUB ebooks."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        texts = []
        with zipfile.ZipFile(str(path)) as zf:
            for name in zf.namelist():
                if name.endswith((".xhtml", ".html", ".htm", ".xml")):
                    try:
                        content = zf.read(name).decode(errors="replace")
                        root = ET.fromstring(content)
                        for elem in root.iter():
                            if elem.text and elem.text.strip():
                                texts.append(elem.text.strip())
                    except Exception:
                        continue
        return "\n".join(texts)[:50000]
    except Exception:
        return ""


def _extract_ipynb(path: Path) -> str:
    """Extract text from Jupyter notebooks."""
    try:
        nb = json.loads(path.read_text())
        texts = []
        for cell in nb.get("cells", []):
            source = "".join(cell.get("source", []))
            if source.strip():
                cell_type = cell.get("cell_type", "code")
                texts.append(f"[{cell_type}]\n{source}")
        return "\n\n".join(texts)[:50000]
    except Exception:
        return ""


def _extract_archive_listing(path: Path) -> str:
    """List contents of archive files."""
    try:
        import zipfile
        if zipfile.is_zipfile(str(path)):
            with zipfile.ZipFile(str(path)) as zf:
                names = zf.namelist()[:200]
                return f"[Archive: {len(names)} files]\n" + "\n".join(names)
    except Exception:
        pass
    return f"[Archive file] {path.name}"


def _extract_media_metadata(path: Path) -> str:
    """Extract metadata from audio/video files."""
    try:
        import subprocess
        result = subprocess.run(
            ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", str(path)],
            capture_output=True, text=True, timeout=10
        )
        if result.returncode == 0:
            data = json.loads(result.stdout)
            fmt = data.get("format", {})
            tags = fmt.get("tags", {})
            parts = []
            for key in ("title", "artist", "album", "genre", "date", "comment", "description"):
                val = tags.get(key) or tags.get(key.upper())
                if val:
                    parts.append(f"{key}: {val}")
            duration = fmt.get("duration")
            if duration:
                parts.append(f"duration: {float(duration):.0f}s")
            return "\n".join(parts) if parts else f"[Media file] {path.name}"
    except Exception:
        pass
    return f"[Media file] {path.name}"


# --- AI Classifier ---
#
# Every file is read and classified by AI. No hardcoded rules.
# Cost optimization through content caching only — if identical
# content was classified before, reuse the result.

CLASSIFY_PROMPT = """You are a universal document classifier. Analyze the document content and return JSON.

The user can be ANYONE — a student, engineer, lawyer, doctor, accountant, business owner, etc.
Do NOT assume any specific profession, country, or industry.
Classify based purely on what the document CONTAINS.

{
  "title": "Clean, descriptive title based on document content",
  "category": "Best fit from: Finance, Legal, Medical, Education, Career, Business, Engineering, Technology, Personal, Government, Real Estate, Insurance, Travel, Other",
  "subcategory": "Specific document type (e.g. Tax Return, Invoice, Lease, Prescription, Resume, Patent, etc.)",
  "tags": ["relevant", "searchable", "tags"],
  "date": "YYYY-MM-DD if any date found in document, else null",
  "language": "ISO 639-1 code (en, es, fr, de, zh, hi, ar, etc.)",
  "summary": "1-2 sentence factual summary of the document",
  "people": ["names of people mentioned"],
  "organizations": ["companies, institutions, agencies mentioned"],
  "confidence": 0.0 to 1.0
}

Rules:
- Title must describe the CONTENT, not the filename
- Category should be the broadest correct fit
- Subcategory should be the specific document type
- Tags should help someone find this document later
- Extract ALL dates, people, and organizations you can find
- Set confidence based on how clear the document content is
- Return ONLY valid JSON, no other text"""

# Content cache: hash of extracted text -> classification result
_content_cache: dict[str, dict] = {}


def _content_cache_key(text: str) -> str:
    return hashlib.md5(text[:2000].encode()).hexdigest()


def classify_document(text: str, filename: str, mime_type: str, file_path: str = "") -> dict:
    """Classify by reading content. Cache identical content to save API calls."""

    # Cache hit — identical content already classified
    if text and len(text) > 20:
        cache_key = _content_cache_key(text)
        if cache_key in _content_cache:
            cached = _content_cache[cache_key].copy()
            return cached

    # No API key or empty content — basic fallback
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if not api_key or len((text or "").strip()) < 10:
        return _fallback_classify(filename, text or "")

    # Claude classifies based on actual content
    try:
        client = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=512,
            messages=[{
                "role": "user",
                "content": f"{CLASSIFY_PROMPT}\n\nDocument content:\n{text[:4000]}",
            }],
        )
        result_text = response.content[0].text.strip()
        if result_text.startswith("```"):
            result_text = result_text.split("\n", 1)[1].rsplit("```", 1)[0]
        result = json.loads(result_text)

        # Cache for identical content
        if text and len(text) > 20:
            _content_cache[_content_cache_key(text)] = result

        return result
    except Exception:
        return _fallback_classify(filename, text or "")


def _fallback_classify(filename: str, text: str = "") -> dict:
    stem = Path(filename).stem
    ext = Path(filename).suffix.lower()

    AUDIO_EXTS = {".mp3", ".m4a", ".wav", ".flac", ".ogg", ".aac", ".wma", ".opus"}
    VIDEO_EXTS = {".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm", ".m4v", ".mpg", ".mpeg", ".3gp"}

    # Parse embedded metadata tags from ffprobe text (media files)
    meta: dict[str, str] = {}
    if text and text.startswith(("[Media", "title:", "artist:", "duration:")):
        for line in text.splitlines():
            if ": " in line:
                k, _, v = line.partition(": ")
                meta[k.strip().lower()] = v.strip()

    if ext in AUDIO_EXTS:
        # Priority: embedded title tag → "Artist - Title" → cleaned filename
        if meta.get("title") and meta.get("artist"):
            title = f"{meta['artist']} — {meta['title']}"
        elif meta.get("title"):
            title = meta["title"]
        else:
            # Clean filename: strip track numbers, underscores, dashes
            import re
            title = re.sub(r"^\d+[\s._-]+", "", stem)  # leading track number
            title = title.replace("_", " ").replace(".", " ")
            # "Artist - Title" filename pattern
            if " - " in title:
                parts = title.split(" - ", 1)
                title = f"{parts[0].strip()} — {parts[1].strip()}"
        tags = []
        if meta.get("genre"): tags.append(meta["genre"])
        if meta.get("album"): tags.append(meta["album"])
        return {
            "title": title or stem or "Untitled Audio",
            "category": "Personal",
            "subcategory": "Audio",
            "tags": tags,
            "date": meta.get("date"),
            "language": "en",
            "summary": f"Audio: {title or stem}" + (f" by {meta['artist']}" if meta.get("artist") else ""),
            "people": [meta["artist"]] if meta.get("artist") else [],
            "organizations": [],
            "confidence": 0.3,
        }

    if ext in VIDEO_EXTS:
        if meta.get("title"):
            title = meta["title"]
        else:
            import re
            title = re.sub(r"^\d+[\s._-]+", "", stem)
            title = title.replace("_", " ").replace(".", " ").replace("-", " - ")
            # "S01E02 Title" pattern
            title = re.sub(r"[Ss](\d+)[Ee](\d+)[\s._-]*", r"S\1E\2 ", title).strip()
        return {
            "title": title or stem or "Untitled Video",
            "category": "Personal",
            "subcategory": "Video",
            "tags": [],
            "date": meta.get("date"),
            "language": "en",
            "summary": f"Video: {title or stem}",
            "people": [],
            "organizations": [],
            "confidence": 0.3,
        }

    name = stem.replace("_", " ").replace("-", " ")
    return {
        "title": name or "Untitled",
        "category": "Other",
        "subcategory": "Uncategorized",
        "tags": [],
        "date": None,
        "language": "en",
        "summary": f"File: {filename}",
        "people": [],
        "organizations": [],
        "confidence": 0.1,
    }


# --- Search ---


def get_search_client():
    try:
        return meilisearch.Client(MEILI_URL, MEILI_KEY)
    except Exception:
        return None


def init_search():
    client = get_search_client()
    if not client:
        print("Meilisearch not available — using SQLite fallback search")
        return
    try:
        client.create_index("documents", {"primaryKey": "id"})
    except Exception:
        pass
    try:
        idx = client.index("documents")
        idx.update_searchable_attributes(["title", "summary", "full_text", "tags", "people", "organizations", "category"])
        idx.update_filterable_attributes(["category", "subcategory", "tags", "language", "is_archived"])
        idx.update_sortable_attributes(["created_at", "document_date", "title", "file_size"])
    except Exception:
        print("Meilisearch config failed — using SQLite fallback search")


def index_document(doc: dict):
    client = get_search_client()
    if not client:
        return
    try:
        search_doc = {k: v for k, v in doc.items() if k != "full_text"}
        search_doc["full_text"] = (doc.get("full_text") or "")[:50000]
        client.index("documents").add_documents([search_doc])
    except Exception:
        pass  # Meilisearch offline — SQLite fallback handles search


# --- Core: Scan & Index ---


def compute_hash(file_path: str) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def scan_and_index(file_path: str, source: str = "scan") -> dict | None:
    """Index a single file in-place. No copy, no upload."""
    path = Path(file_path)
    if not path.exists() or not path.is_file():
        return None
    if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        return None
    if path.name.startswith(".") or path.name.startswith("~$"):
        return None

    conn = get_db()
    try:
        existing = conn.execute(
            "SELECT id FROM documents WHERE file_path = ? AND is_deleted = 0",
            (str(path),)
        ).fetchone()
        if existing:
            return None  # already indexed

        file_hash = compute_hash(str(path))
        hash_dupe = conn.execute(
            "SELECT id, title FROM documents WHERE file_hash = ? AND is_deleted = 0",
            (file_hash,)
        ).fetchone()

        mime_type = magic.from_file(str(path), mime=True)
        file_size = path.stat().st_size
        full_text = extract_text(str(path), mime_type)
        classification = classify_document(full_text, path.name, mime_type, str(path))

        doc_id = str(uuid.uuid4())
        now = datetime.now(timezone.utc).isoformat()

        conn.execute("""
            INSERT INTO documents (
                id, file_path, title, original_filename, category, subcategory,
                file_hash, mime_type, file_size, language, summary, full_text,
                ai_confidence, document_date, tags, people, organizations,
                ai_metadata, source, created_at, indexed_at
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            doc_id, str(path), classification.get("title", path.stem), path.name,
            classification.get("category"), classification.get("subcategory"),
            file_hash, mime_type, file_size,
            classification.get("language", "en"), classification.get("summary"),
            full_text, classification.get("confidence"),
            classification.get("date"),
            json.dumps(classification.get("tags", [])),
            json.dumps(classification.get("people", [])),
            json.dumps(classification.get("organizations", [])),
            json.dumps(classification), source, now, now,
        ))
        conn.commit()

        doc_dict = {
            "id": doc_id, "title": classification.get("title", path.stem),
            "file_path": str(path), "category": classification.get("category"),
            "subcategory": classification.get("subcategory"),
            "summary": classification.get("summary"),
            "full_text": full_text, "tags": classification.get("tags", []),
            "people": classification.get("people", []),
            "organizations": classification.get("organizations", []),
            "file_size": file_size, "mime_type": mime_type,
            "document_date": classification.get("date"),
            "created_at": now, "is_archived": False,
            "is_duplicate": hash_dupe is not None,
        }
        index_document(doc_dict)
        return doc_dict

    finally:
        conn.close()


def scan_directory(directory: str, progress_callback=None) -> dict:
    """Walk a directory tree, index all supported files."""
    root = Path(directory)
    if not root.exists():
        return {"error": f"Directory not found: {directory}"}

    files = [
        p for p in root.rglob("*")
        if p.is_file()
        and p.suffix.lower() in SUPPORTED_EXTENSIONS
        and not p.name.startswith(".")
        and not p.name.startswith("~$")
    ]

    total = len(files)
    results = {"total": total, "indexed": 0, "skipped": 0, "failed": 0, "errors": []}

    for i, path in enumerate(files):
        try:
            result = scan_and_index(str(path), source="scan")
            if result:
                results["indexed"] += 1
            else:
                results["skipped"] += 1
        except Exception as e:
            results["failed"] += 1
            results["errors"].append({"file": path.name, "error": str(e)})

        if progress_callback and (i + 1) % 5 == 0:
            progress_callback(i + 1, total, results)

    return results


# --- FastAPI App ---

app = FastAPI(title="DocuVault Desktop", version="2.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

# Register all feature endpoints
try:
    from features import router as features_router
    app.include_router(features_router)
except ImportError:
    pass

try:
    from features_extra import router as extra_router
    app.include_router(extra_router)
except ImportError:
    pass

try:
    from features_views import router as views_router
    app.include_router(views_router)
except ImportError:
    pass

try:
    from features_collab import router as collab_router
    app.include_router(collab_router)
except ImportError:
    pass

try:
    from features_ai import router as ai_router
    app.include_router(ai_router)
except ImportError:
    pass

try:
    from features_platform import router as platform_router
    app.include_router(platform_router)
except ImportError:
    pass


@app.on_event("startup")
def startup():
    init_db()
    init_search()


@app.get("/api/health")
def health():
    return {"status": "ok", "db": str(DB_PATH)}


@app.get("/api/stats")
def stats():
    conn = get_db()
    total = conn.execute("SELECT COUNT(*) FROM documents WHERE is_deleted = 0").fetchone()[0]
    size = conn.execute("SELECT COALESCE(SUM(file_size), 0) FROM documents WHERE is_deleted = 0").fetchone()[0]
    cats = dict(conn.execute(
        "SELECT COALESCE(category, 'Other'), COUNT(*) FROM documents WHERE is_deleted = 0 GROUP BY category"
    ).fetchall())
    pending = conn.execute("SELECT COUNT(*) FROM documents WHERE is_deleted = 0 AND ai_confidence < 0.5").fetchone()[0]
    folders = conn.execute("SELECT COUNT(*) FROM watched_folders WHERE is_active = 1").fetchone()[0]
    # Classification cost tracking
    rule_based = conn.execute(
        "SELECT COUNT(*) FROM documents WHERE is_deleted = 0 AND ai_confidence <= 0.7 AND ai_confidence >= 0.5"
    ).fetchone()[0]
    ai_classified = conn.execute(
        "SELECT COUNT(*) FROM documents WHERE is_deleted = 0 AND ai_confidence > 0.7"
    ).fetchone()[0]
    api_cost_estimate = ai_classified * 0.005  # ~$0.005 per Haiku call

    conn.close()
    return {
        "total_documents": total, "total_size_bytes": size, "categories": cats,
        "pending_review": pending, "watched_folders": folders,
        "classification": {
            "rule_based_free": total - ai_classified,
            "ai_classified": ai_classified,
            "estimated_api_cost": f"${api_cost_estimate:.2f}",
            "cached_patterns": len(_content_cache),
        },
    }


@app.get("/api/documents")
def list_documents(category: str = None, is_archived: int = 0, offset: int = 0, limit: int = 50):
    conn = get_db()
    query = "SELECT * FROM documents WHERE is_deleted = 0 AND is_archived = ?"
    params = [is_archived]
    if category:
        query += " AND category = ?"
        params.append(category)
    query += " ORDER BY indexed_at DESC LIMIT ? OFFSET ?"
    params.extend([limit, offset])

    docs = [_row_to_dict(r) for r in conn.execute(query, params).fetchall()]
    total = conn.execute(
        "SELECT COUNT(*) FROM documents WHERE is_deleted = 0 AND is_archived = ?" + (" AND category = ?" if category else ""),
        [is_archived] + ([category] if category else [])
    ).fetchone()[0]
    conn.close()
    return {"documents": docs, "total": total, "offset": offset, "limit": limit}


@app.get("/api/documents/{doc_id}")
def get_document(doc_id: str):
    conn = get_db()
    row = conn.execute("SELECT * FROM documents WHERE id = ? AND is_deleted = 0", (doc_id,)).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Not found"}, 404)
    return _row_to_dict(row)


@app.get("/api/documents/{doc_id}/open")
def open_document(doc_id: str):
    """Open file in default system app."""
    conn = get_db()
    row = conn.execute("SELECT file_path FROM documents WHERE id = ?", (doc_id,)).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Not found"}, 404)
    file_path = row["file_path"]
    if not Path(file_path).exists():
        return JSONResponse({"error": "File not found on disk"}, 404)

    import subprocess, platform
    system = platform.system()
    if system == "Darwin":
        subprocess.Popen(["open", file_path])
    elif system == "Windows":
        os.startfile(file_path)
    else:
        subprocess.Popen(["xdg-open", file_path])
    return {"status": "opened", "path": file_path}


@app.get("/api/documents/{doc_id}/reveal")
def reveal_in_finder(doc_id: str):
    """Show file in Finder/Explorer."""
    conn = get_db()
    row = conn.execute("SELECT file_path FROM documents WHERE id = ?", (doc_id,)).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Not found"}, 404)
    file_path = row["file_path"]

    import subprocess, platform
    system = platform.system()
    if system == "Darwin":
        subprocess.Popen(["open", "-R", file_path])
    elif system == "Windows":
        subprocess.Popen(["explorer", "/select,", file_path])
    return {"status": "revealed", "path": file_path}


@app.put("/api/documents/{doc_id}")
def update_document(doc_id: str, data: dict):
    conn = get_db()
    fields = []
    params = []
    for key in ("title", "category", "subcategory", "is_archived"):
        if key in data:
            fields.append(f"{key} = ?")
            params.append(data[key])
    if "tags" in data:
        fields.append("tags = ?")
        params.append(json.dumps(data["tags"]))
    if not fields:
        return {"error": "No fields to update"}
    params.append(doc_id)
    conn.execute(f"UPDATE documents SET {', '.join(fields)} WHERE id = ?", params)
    conn.commit()
    conn.close()
    return {"status": "updated"}


@app.delete("/api/documents/{doc_id}")
def delete_document(doc_id: str):
    conn = get_db()
    conn.execute("UPDATE documents SET is_deleted = 1 WHERE id = ?", (doc_id,))
    conn.commit()
    conn.close()
    return {"status": "deleted"}


@app.post("/api/documents/{doc_id}/reclassify")
def reclassify(doc_id: str):
    conn = get_db()
    row = conn.execute("SELECT * FROM documents WHERE id = ?", (doc_id,)).fetchone()
    if not row:
        conn.close()
        return JSONResponse({"error": "Not found"}, 404)
    text = row["full_text"] or ""
    classification = classify_document(text, row["original_filename"] or "", row["mime_type"] or "")
    conn.execute("""
        UPDATE documents SET title=?, category=?, subcategory=?, summary=?,
        ai_confidence=?, tags=?, people=?, organizations=?, ai_metadata=?,
        document_date=? WHERE id=?
    """, (
        classification.get("title"), classification.get("category"),
        classification.get("subcategory"), classification.get("summary"),
        classification.get("confidence"),
        json.dumps(classification.get("tags", [])),
        json.dumps(classification.get("people", [])),
        json.dumps(classification.get("organizations", [])),
        json.dumps(classification), classification.get("date"), doc_id,
    ))
    conn.commit()
    conn.close()
    return classification


@app.get("/api/search")
def search(q: str = "", category: str = None, offset: int = 0, limit: int = 20):
    client = get_search_client()
    if not client or not q:
        # Fallback to SQLite FTS
        conn = get_db()
        like = f"%{q}%"
        query = "SELECT * FROM documents WHERE is_deleted = 0 AND (title LIKE ? OR summary LIKE ? OR full_text LIKE ?)"
        params = [like, like, like]
        if category:
            query += " AND category = ?"
            params.append(category)
        query += " LIMIT ? OFFSET ?"
        params.extend([limit, offset])
        docs = [_row_to_dict(r) for r in conn.execute(query, params).fetchall()]
        conn.close()
        return {"hits": docs, "total": len(docs), "query": q, "processing_time_ms": 0}

    try:
        filters = f'category = "{category}"' if category else None
        params = {"offset": offset, "limit": limit}
        if filters:
            params["filter"] = filters
        result = client.index("documents").search(q, params)
        return {
            "hits": result.get("hits", []),
            "total": result.get("estimatedTotalHits", 0),
            "query": q,
            "processing_time_ms": result.get("processingTimeMs", 0),
        }
    except Exception:
        # Meilisearch unavailable — fallback to SQLite
        conn = get_db()
        like = f"%{q}%"
        query = "SELECT * FROM documents WHERE is_deleted = 0 AND (title LIKE ? OR summary LIKE ? OR full_text LIKE ?)"
        params = [like, like, like]
        if category:
            query += " AND category = ?"
            params.append(category)
        query += " LIMIT ? OFFSET ?"
        params.extend([limit, offset])
        docs = [_row_to_dict(r) for r in conn.execute(query, params).fetchall()]
        conn.close()
        return {"hits": docs, "total": len(docs), "query": q, "processing_time_ms": 0}


@app.get("/api/facets")
def facets():
    conn = get_db()
    cats = dict(conn.execute(
        "SELECT COALESCE(category, 'Other'), COUNT(*) FROM documents WHERE is_deleted = 0 GROUP BY category"
    ).fetchall())
    conn.close()
    return {"category": cats}


@app.get("/api/duplicates")
def duplicates():
    conn = get_db()
    dupes = conn.execute("""
        SELECT file_hash, COUNT(*) as cnt FROM documents
        WHERE is_deleted = 0 AND file_hash IS NOT NULL
        GROUP BY file_hash HAVING cnt > 1
    """).fetchall()

    result = []
    for row in dupes:
        docs = conn.execute(
            "SELECT id, title, file_path, file_size, indexed_at, category FROM documents WHERE file_hash = ? AND is_deleted = 0",
            (row["file_hash"],)
        ).fetchall()
        result.append({
            "file_hash": row["file_hash"],
            "count": row["cnt"],
            "documents": [dict(d) for d in docs],
        })
    conn.close()
    return {"hash_duplicates": result}


@app.post("/api/scan")
def start_scan(directory: str, dry_run: bool = True):
    """Scan a directory and index all files."""
    if dry_run:
        root = Path(directory)
        if not root.exists():
            return JSONResponse({"error": f"Not found: {directory}"}, 404)
        files = [p for p in root.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS]
        return {"total": len(files), "directory": directory, "dry_run": True}

    result = scan_directory(directory)

    conn = get_db()
    folder_id = str(uuid.uuid4())
    conn.execute(
        "INSERT OR REPLACE INTO watched_folders (id, path, last_scan, file_count, is_active) VALUES (?, ?, ?, ?, 1)",
        (folder_id, directory, datetime.now(timezone.utc).isoformat(), result["indexed"])
    )
    conn.commit()
    conn.close()
    return result


@app.get("/api/watched-folders")
def list_watched_folders():
    conn = get_db()
    folders = [dict(r) for r in conn.execute("SELECT * FROM watched_folders WHERE is_active = 1").fetchall()]
    conn.close()
    return folders


@app.get("/api/export")
def export_metadata():
    conn = get_db()
    docs = [_row_to_dict(r) for r in conn.execute("SELECT * FROM documents WHERE is_deleted = 0").fetchall()]
    conn.close()
    return {"exported_at": datetime.now(timezone.utc).isoformat(), "total": len(docs), "documents": docs}


def _row_to_dict(row) -> dict:
    d = dict(row)
    for field in ("tags", "people", "organizations"):
        if isinstance(d.get(field), str):
            try:
                d[field] = json.loads(d[field])
            except (json.JSONDecodeError, TypeError):
                d[field] = []
    return d


# --- Chat & Intelligence ---

_chat_history: list[dict] = []


@app.post("/api/chat")
def chat_endpoint(data: dict):
    """Ask questions about your documents. Conversational AI."""
    from chat import chat
    global _chat_history

    message = data.get("message", "")
    if not message:
        return JSONResponse({"error": "No message"}, 400)

    # Add user message to history
    _chat_history.append({"role": "user", "content": message})

    result = chat(message, _chat_history)

    # Add assistant response to history
    _chat_history.append({"role": "assistant", "content": result["response"]})

    # Keep history manageable
    if len(_chat_history) > 20:
        _chat_history = _chat_history[-20:]

    return result


@app.delete("/api/chat/history")
def clear_chat():
    """Clear conversation history."""
    global _chat_history
    _chat_history = []
    return {"status": "cleared"}


@app.get("/api/insights")
def get_insights():
    """Proactive alerts: expiring docs, duplicates, review needed."""
    from chat import proactive_insights
    return proactive_insights()


if __name__ == "__main__":
    port = int(sys.argv[1]) if len(sys.argv) > 1 else 8200
    uvicorn.run(app, host="127.0.0.1", port=port, log_level="info")
